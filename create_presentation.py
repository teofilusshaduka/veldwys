import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

def add_slide(title_text, content_text, notes_text, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    if content_text is not None and len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = content_text
        
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text
    return slide

# Slide 0: Title and framing
add_slide(
    "VeldWys: Farm Command Centre",
    "Beyond a Chatbot: A comprehensive farm management system for Namibian farmers.",
    "We got tasked in this project to create an AI chatbot. And we thought, hey, that's cool, but why don't we take it a step further and actually make a whole farm management system which actually does way more for the farmer than a simple chatbot would.",
    layout_index=0
)

# Slide 1
add_slide(
    "The farmer and the notebook",
    "[Photo of the real handwritten stock book: Ear tags, colours, Ewe / Ram / Kapater]",
    "Livestock is the backbone of the rural Namibian economy, and for most farmers the herd isn't just an asset — it's the savings account, the pension, and the inheritance, all at once. What you're looking at is a real page from a real farmer's notebook. Twenty years of records like this exist across the country: the entire asset register of a farm, written in pen, kept in one notebook, in one place, with no backup anywhere.\n\nEvery piece of digital farming software offered to this farmer so far has assumed he'd sit down and retype all of it. He won't do that, and honestly, he shouldn't have to.\n\nWe didn't start by asking what AI could do for farmers. We started with the notebook."
)

# Slide 2
add_slide(
    "Register first, chatbot second",
    "Generic farming chatbot:\n\"How many cattle do you have? Where is your farm? What breed?\" (every single time)\n\nVeldWys:\nKnows TE-TANGA S009 is a brown ewe with a white spot on her tail, vaccinated in March.",
    "Anyone can put a grazing prompt in front of a language model and call it a farming assistant. That's a demo, not a product, because it forgets the farm the moment you close the tab. VeldWys is built the other way around: it's a livestock register first, with an advisor sitting on top of it. Every question the farmer asks carries the real context of that farm automatically — herd composition, ear tags, camp size, region, rainfall, vaccination history. The farmer never has to restate any of it. That's the entire design philosophy in one sentence.\n\nAnd it isn't read-only. The advisor has nine tools it can actually use, including writing back to the register — so if a farmer says \"I sold three goats,\" three animal records actually flip to sold, right there in the conversation.\n\nThe AI is the thin part. The register is the product."
)

# Slide 3
add_slide(
    "DEMO — the notebook scan",
    "[Notebook photo on the left] -> [Extracted editable records on the right, with a Kapater row and an untagged row visibly flagged]",
    "This is the feature that decides whether the app actually gets used or not. If onboarding means manual typing, nobody onboards — full stop. So we built a scan that turns a photograph of a handwritten page directly into editable records.\n\nThe important part is that it isn't tuned to my notebook specifically. The pipeline reads the page on its own terms first — is this even a table, what does each column seem to hold, are there section headers grouping the rows — and only after that does it map what it found onto our schema. That means it copes with column order changing, headers being missing entirely, free-text lists, a vaccination diary, even a tally sheet, because it isn't pattern-matching against one layout.\n\nIt also handles the messy reality of a real notebook: ear tags written across two lines, rows with no tag at all, rows that have been struck through and must not be imported, duplicate tags, even a phone number scribbled in the margin that shouldn't end up in a field. And \"Kapater\" — a castrated male — had no home at all in a simple male-or-female schema. Real registers need that distinction, so now ours has it.\n\nNothing gets saved until the farmer confirms it on screen. Uncertain records float to the top, and the app tells you plainly how it read the page, so the farmer is always in control of what lands in their register.\n\n(If a judge asks \"did you train a model?\"): No, and we didn't need to. The gain came from splitting the job into two steps and never showing the model the target format until the second one."
)

# Slide 4
add_slide(
    "DEMO — voice, in the farmer's language",
    "[Full-screen voice view, orb mid-pulse]",
    "Text input quietly assumes two things: that you're comfortable reading and typing, and that you can do it in your own language. Neither is a safe assumption here, and neither is realistic for a farmer standing in a kraal with gloves on. So VeldWys has a full, hands-free conversation mode. It listens, hears when you've paused, answers out loud, and starts listening again — no tapping the screen at any point.\n\nIt also adapts to the room it's in. At the start of every turn it quietly samples the background noise and sets its listening threshold from that, so the same app works in a quiet kitchen and next to a running bakkie without being retuned. And it knows not to transcribe its own spoken reply back as the next question.\n\nAll of this works end to end in five languages — the question, the answer, and the spoken reply.\n\nOshiwambo speech recognition specifically is still weak, because the underlying model has almost no training data for it. We fixed what was fixable — it no longer silently answers in English when it mishears — and we're upfront about what isn't fixed yet. Oshiwambo as written text, and Oshiwambo spoken back to the farmer, are both solid."
)

# Slide 5
add_slide(
    "Getting the language right",
    "Was -> Literally means -> Now\n\n• omugongo gwoshimeni -> \"the spine of anthrax\" -> ondjeka yaanthrax\n• dhi na omakutsi -> \"they have ears\" -> dhi na iidhindilo\n• iikadhona -> \"girls\" (for cows) -> oonkadhi\n• Omulumentu / Omukiintu -> human man / woman -> Ondume / Onkadhi",
    "We shipped with four languages in our first round. Then a native Oshiwambo speaker opened the app and told us it simply didn't read like Oshiwambo. He was right, and it wasn't a polish problem — it was three real faults at once. The file mixed two different orthographies in the same document, the grammar around noun classes was broken throughout, and in several places we had, honestly, just the wrong word.\n\nOnce we fixed that, we split the languages up properly. Oshindonga and Oshikwanyama had been collapsed into a single language when they're not the same, and Otjiherero was sitting on Oshikwanyama's own ISO code. Today we have five languages, 286 strings each, with full parity across all of them.\n\nWe also built a review tool so a native speaker can correct the file directly, rather than reporting a bug to a developer who then has to guess again at the right word.\n\nAnyone can run text through a translation API. Knowing that it says \"the spine of anthrax\" instead of \"anthrax injection\" takes an actual speaker in the room. We built the workflow that keeps one there."
)

# Slide 6
add_slide(
    "What it costs to run",
    "Measured cost per operation:\n• Chat question: $0.00064\n• Voice question: $0.00068\n• Spoken answer: $0.0021\n• Full voice turn: $0.0034\n• Notebook page scan: $0.056\n\nPer farmer:\n• Per month (60 text + 20 voice): $0.11 / N$2.00\n• One-time digitising (20-page notebook): $1.12 / N$21\n• Year one, all in: $2.40 / N$44",
    "A full year of AI advice for one farmer costs less than a single cup of coffee. Digitising twenty years of paper records costs about N$21, and that's a one-time cost, not a recurring one.\n\nI want to be clear that this isn't an estimate we made up to sound good. We instrumented the running application and read the actual token counts off real requests — the method is documented at the end of this deck if anyone wants to check our working.\n\nWe got these numbers by being deliberate about where the money actually goes. Two things in this app cost real money: reading handwriting, and reasoning through a question, so those are the only places we spend on strong models. The vaccination schedules and the proactive insight rules are plain Python with zero LLM cost, because a calendar simply doesn't need a language model to run it. And the one genuinely expensive operation — scanning a notebook — is something that happens once per farmer, ever. The thing that happens every single day costs about six hundredths of a US cent.\n\nAt N$2 a month, this scales to every communal farmer in Namibia, and the AI bill is not what's going to stop us."
)

# Slide 7
add_slide(
    "How we know it works",
    "89/89 checks passed across 6 page formats",
    "We didn't just eyeball the scan and call it done — we graded it against notebook pages laid out six genuinely different ways: no header row at all, columns in a different order, a free-text list instead of a table, a date-led diary, a tally sheet, and an Afrikaans-labelled form.\n\nThe two hardest cases in that set are the ones that matter most. A tally sheet that just says \"cattle 47, goats 112\" has to produce two counts, not 226 invented animal records. And a row that's been struck through on the page has to be recognised as deleted and left out entirely. It passes both of those, along with everything else — eighty-nine out of eighty-nine checks.\n\nWe also deliberately checked that our own test suite is capable of failing, because a green suite that can never turn red doesn't actually prove anything.\n\nSo when we say it works on notebooks it has never seen before, that isn't a hope — it's something we specifically measured."
)

# Slide 8
add_slide(
    "What's real, and what's next",
    "Working today:\n• Livestock register\n• Notebook scan\n• 5 languages\n• Voice conversation\n• Offline-first PWA\n• Namibian vaccination schedules\n• Proactive insight engine\n• Real rangeland data\n• Analytics\n• Password recovery\n\nKnown gaps:\n• Oshiwambo ASR is weak\n• TTS accents need native tuning\n• API endpoints need auth\n• No SMS for password recovery",
    "We'd rather name our own weak spots than have a judge find them first. Every one of these gaps is going to be visible to someone in this room eventually, so we're putting them on the slide ourselves — being first with that is worth more to us than being lucky. And to be clear, none of these gaps are architectural problems. They're all scoped, understood, and on the list."
)

# Slide 9
add_slide(
    "Close",
    "VeldWys: The notebook was never the problem. Nobody had bothered to read it.",
    "Namibian farmers have been keeping excellent records for decades. They never needed to be taught how to collect data — what they needed was a piece of software willing to meet them where that data already lives, which is a paper notebook, not a form.\n\nVeldWys reads that notebook, speaks the farmer's own language, works with no signal at all, and costs about N$2 per farmer per month to run.\n\nThe notebook was never the problem. Nobody had bothered to read it."
)

# Slide 10: Appendix
add_slide(
    "Appendix: Cost Methodology",
    "Measured token counts off real API responses.\n• Opus 5: $5 in / $25 out\n• Sonnet 5: $3 in / $15 out\n• GPT-4o-mini: $0.15 in / $0.60 out\n• Whisper: $0.006 / min\n\n18.5 N$ = $1 USD",
    "Every figure comes from instrumenting the running application and reading the actual usage token counts off real API responses, not from a pricing calculator or a rough guess. The chat cost is averaged over four real farm questions asked against a seeded herd, including turns where the agent calls tools and makes two API round trips instead of one. The scan cost is measured on a real fixture page, with both passes and their input and output tokens counted separately.\n\nThe published prices we used, per million tokens, are: Opus 5 at $5 input and $25 output, Sonnet 5 at $3 and $15, GPT-4o-mini at $0.15 and $0.60, and Whisper at $0.006 per minute of audio. Sonnet 5 is actually on introductory pricing right now at $2 and $10, which would make the scan cost roughly 10% cheaper than what we're showing — we used the full list price on purpose, so our number is the conservative one, not the flattering one.\n\nThe monthly model assumes sixty text questions and twenty voice turns per farmer per month; that assumption can be scaled however a judge prefers, since the real output here is the per-operation cost, not the monthly multiplier. All rand figures use a conversion rate of 18.5 N$ to the US dollar."
)

prs.save("VeldWys_Presentation.pptx")
print("Presentation generated successfully!")
